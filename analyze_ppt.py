# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

path = r'c:\Users\66496\.trae-cn\attachments\6a7ec78c2d0cfdb32a600535\a4960b04-6656-4941-8dad-52b81a3908f8_8743cdee-25ad-4054-ab14-13371a8358f6_营员路演PPT.pptx'
prs = Presentation(path)

for i, slide in enumerate(prs.slides):
    print(f'===== Slide {i+1} =====')
    for shape in slide.shapes:
        print(f'  Shape: {shape.name}, type={shape.shape_type}')
        print(f'    pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})')
        try:
            fill = shape.fill
            if fill.type is not None:
                try:
                    print(f'    Fill color: {fill.fore_color.rgb}')
                except:
                    print(f'    Fill type: {fill.type}')
        except:
            pass
        if shape.has_text_frame:
            for j, para in enumerate(shape.text_frame.paragraphs):
                t = para.text.strip()
                if t:
                    print(f'    Para {j}: {t[:100]}')
                    print(f'      align={para.alignment}')
                    for run in para.runs:
                        fc = ''
                        try:
                            fc = f'color={run.font.color.rgb}'
                        except:
                            fc = 'color=inherited'
                        print(f'      Run: font={run.font.name} size={run.font.size} bold={run.font.bold} {fc}')
    print()

# -*- coding: utf-8 -*-
"""青囊AI 项目演示PPT生成脚本 — 四大板块版"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 色彩体系
# ============================================================
JADE       = RGBColor(0x0D, 0x6B, 0x5B)
JADE_DEEP  = RGBColor(0x07, 0x3B, 0x35)
JADE_LIGHT = RGBColor(0xE7, 0xF4, 0xEF)
GOLD       = RGBColor(0xC8, 0x9B, 0x47)
GOLD_LIGHT = RGBColor(0xF5, 0xED, 0xD6)
PAPER      = RGBColor(0xFA, 0xF8, 0xF3)
INK        = RGBColor(0x17, 0x33, 0x2F)
VERMILION  = RGBColor(0xB5, 0x64, 0x58)
MUTED      = RGBColor(0x6A, 0x7D, 0x78)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LINE_C     = RGBColor(0xD0, 0xDD, 0xD9)

IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static", "images")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN

# ============================================================
# 工具函数
# ============================================================
def add_bg(slide, color=PAPER):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, x, y, w, h, text, fs=14, color=INK, bold=False, fn="Microsoft YaHei",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.5):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fs)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = fn
    p.alignment = align
    p.space_after = Pt(0)
    if ls:
        p.line_spacing = ls
    return tb

def add_multi(slide, x, y, w, h, lines, fs=14, color=INK, fn="Microsoft YaHei",
              align=PP_ALIGN.LEFT, ls=1.6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fs)
        p.font.color.rgb = color
        p.font.name = fn
        p.alignment = align
        p.line_spacing = ls
        p.space_after = Pt(4)
    return tb

def title_bar(slide, title, subtitle="", section=""):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), fill_color=JADE_DEEP)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Pt(3), fill_color=GOLD)
    add_text(slide, MARGIN, Inches(0.15), Inches(8), Inches(0.55),
             title, fs=28, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, MARGIN, Inches(0.70), Inches(8), Inches(0.35),
                 subtitle, fs=13, color=GOLD_LIGHT)
    if section:
        add_text(slide, SLIDE_W - Inches(3.5), Inches(0.25), Inches(2.8), Inches(0.5),
                 section, fs=11, color=GOLD, align=PP_ALIGN.RIGHT)

def footer(slide, n):
    add_text(slide, MARGIN, SLIDE_H - Inches(0.4), Inches(6), Inches(0.3),
             "青囊AI · 中医体质智能分析平台", fs=9, color=MUTED)
    add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.4), Inches(1), Inches(0.3),
             str(n), fs=9, color=MUTED, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, fill=WHITE, border=LINE_C):
    return add_rect(slide, x, y, w, h, fill_color=fill, line_color=border, line_width=Pt(1))

def icon_circle(slide, x, y, d, txt, bg=JADE, tc=WHITE, fs=20):
    c = add_rect(slide, x, y, d, d, fill_color=bg, shape_type=MSO_SHAPE.OVAL)
    tf = c.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(fs)
    p.font.color.rgb = tc
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return c

def tag_pill(slide, x, y, w, h, text, bg=JADE_LIGHT, tc=JADE, fs=9):
    s = add_rect(slide, x, y, w, h, fill_color=bg, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = s.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fs)
    p.font.color.rgb = tc
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return s


# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, JADE_DEEP)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=JADE_DEEP)
for cx, cy, r in [(Inches(10.8), Inches(1.2), Inches(2.2)), (Inches(1.5), Inches(5.8), Inches(1.6))]:
    add_rect(slide, cx - r, cy - r, r * 2, r * 2, fill_color=JADE, shape_type=MSO_SHAPE.OVAL)

add_rect(slide, Inches(4), Inches(2.1), Inches(5.333), Pt(2), fill_color=GOLD)
add_text(slide, Inches(2), Inches(2.4), Inches(9.333), Inches(1.2),
         "青囊AI", fs=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(3.6), Inches(9.333), Inches(0.6),
         "中医体质智能分析平台", fs=24, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(4.3), Inches(9.333), Inches(0.5),
         "取意「青囊书」——华佗所传中医经典之名", fs=14,
         color=RGBColor(0xAA, 0xBB, 0xB5), align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(4.9), Inches(9.333), Inches(0.5),
         "以AI之眼观舌象 · 以古人之智辨体质", fs=16, color=GOLD, align=PP_ALIGN.CENTER)

tags = ["项目背景", "产品设计过程", "产品展示", "未来展望"]
tw = Inches(2.0)
gap = Inches(0.3)
total = tw * len(tags) + gap * (len(tags) - 1)
sx = (SLIDE_W - total) / 2
for i, t in enumerate(tags):
    tx = sx + i * (tw + gap)
    r = add_rect(slide, tx, Inches(5.7), tw, Inches(0.45),
                 fill_color=None, line_color=GOLD, line_width=Pt(1),
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = r.text_frame
    p = tf.paragraphs[0]
    p.text = t
    p.font.size = Pt(12)
    p.font.color.rgb = GOLD_LIGHT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_text(slide, Inches(2), Inches(6.6), Inches(9.333), Inches(0.4),
         "2026 · 青囊AI项目演示", fs=11, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "目  录", "CONTENTS")

toc = [
    ("01", "项目背景", "市场需求 · 痛点分析 · 项目定位"),
    ("02", "产品设计过程", "需求分析 → 数据制作 → 模型训练 → 开发上线"),
    ("03", "产品展示", "功能介绍 · 商业模式 · 亮点介绍"),
    ("04", "未来展望", "功能规划 · 商业化路径 · 技术升级"),
]
cw = Inches(5.5)
ch = Inches(2.3)
gx = Inches(0.5)
gy = Inches(0.4)
sx = MARGIN
sy = Inches(1.6)

for i, (num, t, d) in enumerate(toc):
    col = i % 2
    row = i // 2
    cx = sx + col * (cw + gx)
    cy = sy + row * (ch + gy)
    card(slide, cx, cy, cw, ch)
    add_rect(slide, cx, cy, Inches(0.08), ch, fill_color=JADE)
    add_text(slide, cx + Inches(0.3), cy + Inches(0.2), Inches(1.5), Inches(0.8),
             num, fs=40, color=GOLD, bold=True, fn="Georgia")
    add_text(slide, cx + Inches(1.8), cy + Inches(0.3), Inches(3.5), Inches(0.5),
             t, fs=20, color=INK, bold=True)
    add_text(slide, cx + Inches(1.8), cy + Inches(0.95), Inches(3.5), Inches(0.8),
             d, fs=12, color=MUTED, ls=1.5)

footer(slide, 2)


# ============================================================
# ===== 第一章：项目背景 =====
# ============================================================

# Slide 3: 项目背景 — 市场需求与痛点
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "市场需求与痛点分析", "为什么需要青囊AI？", "01 项目背景")

# 左：痛点
add_text(slide, MARGIN, Inches(1.5), Inches(5.5), Inches(0.5),
         "当前痛点", fs=18, color=VERMILION, bold=True)
pains = [
    "▸  中医门槛高：普通人难以理解体质概念",
    "    和舌象、面相等诊断指标的含义",
    "",
    "▸  问诊不便：线下中医挂号难、排队久，",
    "    儿童体质调理需求难以满足",
    "",
    "▸  缺乏工具：市面上缺少结合AI技术的",
    "    中医体质科普工具，家长无从下手",
    "",
    "▸  信息碎片化：网络健康信息良莠不齐，",
    "    缺乏系统性体质分析和调理建议",
]
add_multi(slide, MARGIN, Inches(2.1), Inches(5.5), Inches(4.5), pains, fs=12, ls=1.6)

# 右：需求
rx = Inches(7.0)
add_text(slide, rx, Inches(1.5), Inches(5.5), Inches(0.5),
         "用户需求", fs=18, color=JADE, bold=True)
needs = [
    ("🔧", "便捷检测", "在家即可完成舌象/声音/\n面相/问卷四诊检测"),
    ("🤖", "AI分析", "AI自动分析体质倾向，\n生成个性化建议"),
    ("📖", "通俗解读", "用大白话解释体质，\n家长看得懂、用得上"),
    ("🛒", "调理方案", "提供食疗推荐和养生\n产品，形成闭环"),
]
for i, (icon, t, d) in enumerate(needs):
    ny = Inches(2.1) + i * Inches(1.15)
    card(slide, rx, ny, Inches(5.7), Inches(1.0), fill=JADE_LIGHT)
    add_text(slide, rx + Inches(0.2), ny + Inches(0.15), Inches(0.5), Inches(0.6),
             icon, fs=24, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, rx + Inches(0.9), ny + Inches(0.1), Inches(1.8), Inches(0.35),
             t, fs=14, color=JADE, bold=True)
    add_text(slide, rx + Inches(0.9), ny + Inches(0.45), Inches(4.5), Inches(0.5),
             d, fs=11, color=MUTED, ls=1.4)

footer(slide, 3)


# Slide 4: 项目背景 — 项目定位
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "项目定位与理念", "青囊AI是什么？", "01 项目背景")

# 中心定位
add_card = card(slide, MARGIN, Inches(1.5), Inches(12.1), Inches(1.8), fill=JADE_DEEP)
add_text(slide, MARGIN + Inches(0.3), Inches(1.65), Inches(11.5), Inches(0.5),
         "青囊AI · 中医体质智能分析平台", fs=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, MARGIN + Inches(0.3), Inches(2.2), Inches(11.5), Inches(0.8),
         "结合传统中医望闻问切四诊理论，通过YOLO目标检测、DeepSeek AI大模型等技术，\n为普通家庭提供便捷的中医体质分析和调养建议，让千年医道智慧走进每个家庭。",
         fs=13, color=GOLD_LIGHT, align=PP_ALIGN.CENTER, ls=1.6)

# 四大支柱
pillars = [
    ("🌿", "传承", "取意华佗《青囊书》\n传承千年中医智慧", JADE),
    ("🔬", "科技", "YOLO + DeepSeek\nAI赋能体质分析", GOLD),
    ("👨‍👩‍👧", "普惠", "面向普通家庭\n大白话解读体质", VERMILION),
    ("🔄", "闭环", "检测→分析→建议→\n食疗推荐→产品", JADE_DEEP),
]
pw = Inches(2.8)
gap = Inches(0.25)
total = pw * 4 + gap * 3
sx = (SLIDE_W - total) / 2
for i, (icon, t, d, c) in enumerate(pillars):
    px = sx + i * (pw + gap)
    py = Inches(3.7)
    card(slide, px, py, pw, Inches(2.8), fill=WHITE)
    add_rect(slide, px, py, pw, Inches(0.08), fill_color=c)
    add_text(slide, px, py + Inches(0.3), pw, Inches(0.6),
             icon, fs=30, align=PP_ALIGN.CENTER)
    add_text(slide, px, py + Inches(1.0), pw, Inches(0.45),
             t, fs=18, color=c, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, px + Inches(0.2), py + Inches(1.5), pw - Inches(0.4), Inches(1),
             d, fs=11, color=MUTED, align=PP_ALIGN.CENTER, ls=1.5)

footer(slide, 4)


# ============================================================
# ===== 第二章：产品设计过程 =====
# ============================================================

# Slide 5: 产品设计过程 — 总览流程
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "产品设计过程 · 总览", "从需求到上线的完整研发流程", "02 产品设计过程")

steps = [
    ("01", "需求分析", "用户调研\n功能规划", JADE),
    ("02", "数据收集", "舌象图片\n采集整理", GOLD),
    ("03", "数据标注", "21类特征\n标注标注", VERMILION),
    ("04", "数据集制作", "训练集/验证集\n数据增强", JADE_DEEP),
    ("05", "模型训练", "YOLOv8训练\n超参调优", JADE),
    ("06", "模型验证", "mAP评估\n迭代优化", GOLD),
    ("07", "基础网页开发", "Flask框架\n中国风UI", VERMILION),
    ("08", "智能体上线", "青囊Agent\nDeepSeek", JADE_DEEP),
    ("09", "商品广告上线", "虚拟货架\n食疗推荐", JADE),
    ("10", "网站测试", "功能测试\n性能优化", GOLD),
]

cw = Inches(2.3)
ch = Inches(1.5)
gx = Inches(0.12)
gy = Inches(0.2)
cols = 5
sx = (SLIDE_W - (cw * cols + gx * (cols - 1))) / 2
sy = Inches(1.6)

for i, (num, t, d, c) in enumerate(steps):
    col = i % cols
    row = i // cols
    cx = sx + col * (cw + gx)
    cy = sy + row * (ch + gy)
    card(slide, cx, cy, cw, ch, fill=WHITE)
    add_rect(slide, cx, cy, cw, Inches(0.06), fill_color=c)
    add_text(slide, cx + Inches(0.15), cy + Inches(0.15), Inches(0.6), Inches(0.35),
             num, fs=14, color=c, bold=True, fn="Georgia")
    add_text(slide, cx + Inches(0.15), cy + Inches(0.5), cw - Inches(0.3), Inches(0.35),
             t, fs=12, color=INK, bold=True)
    add_text(slide, cx + Inches(0.15), cy + Inches(0.85), cw - Inches(0.3), Inches(0.55),
             d, fs=9, color=MUTED, ls=1.4)
    # 箭头
    if col < cols - 1 and i < len(steps) - 1:
        add_text(slide, cx + cw - Inches(0.05), cy + ch / 2 - Inches(0.1), Inches(0.2), Inches(0.2),
                 "→", fs=12, color=MUTED, align=PP_ALIGN.CENTER)

# 底部说明
add_text(slide, MARGIN, Inches(5.8), Inches(12.1), Inches(0.4),
         "完整研发周期：需求分析 → 数据工程 → 模型训练 → Web开发 → AI功能 → 商业化 → 测试上线",
         fs=12, color=JADE, bold=True, align=PP_ALIGN.CENTER)

footer(slide, 5)


# Slide 6: 需求分析 & 数据收集标注
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "需求分析 · 数据收集 · 数据标注", "明确目标 → 采集数据 → 标注特征", "02 产品设计过程")

# 三栏
col_w = Inches(3.9)
col_gap = Inches(0.2)
sx = MARGIN
col_y = Inches(1.5)
col_h = Inches(5.3)

# 栏1: 需求分析
cx = sx
card(slide, cx, col_y, col_w, col_h)
add_rect(slide, cx, col_y, col_w, Inches(0.5), fill_color=JADE)
add_text(slide, cx + Inches(0.2), col_y + Inches(0.05), col_w, Inches(0.4),
         "01 需求分析", fs=14, color=WHITE, bold=True)
req_lines = [
    "目标用户：",
    "  ▸ 关注儿童体质的家长",
    "  ▸ 对中医养生感兴趣的成年人",
    "",
    "核心功能：",
    "  ▸ 舌象AI检测（望）",
    "  ▸ 声音特征分析（闻）",
    "  ▸ 体质问卷分型（问）",
    "  ▸ 面相分析（切）",
    "  ▸ AI综合健康报告",
    "  ▸ 中医体质科普Agent",
    "",
    "设计原则：",
    "  ▸ 中国风水墨UI",
    "  ▸ 大白话替代术语",
    "  ▸ 页面跳转简洁操作",
]
add_multi(slide, cx + Inches(0.2), col_y + Inches(0.65), col_w - Inches(0.4), col_h - Inches(0.8),
          req_lines, fs=10, ls=1.5)

# 栏2: 数据收集
cx = sx + col_w + col_gap
card(slide, cx, col_y, col_w, col_h)
add_rect(slide, cx, col_y, col_w, Inches(0.5), fill_color=GOLD)
add_text(slide, cx + Inches(0.2), col_y + Inches(0.05), col_w, Inches(0.4),
         "02 数据收集", fs=14, color=WHITE, bold=True)
data_lines = [
    "舌象图片采集：",
    "  ▸ 收集不同舌质颜色的舌象",
    "    （淡红、红、绛、紫、淡白）",
    "  ▸ 收集不同舌苔特征",
    "    （薄白、厚白、薄黄、厚黄、",
    "      腻苔、剥苔、无苔）",
    "  ▸ 收集舌形特征",
    "    （齿痕、胖大、瘦薄、裂纹）",
    "",
    "数据来源：",
    "  ▸ 公开舌象数据集",
    "  ▸ 中医舌诊图谱参考",
    "  ▸ 网络医学教育资源",
    "",
    "数据量：",
    "  ▸ 覆盖21种舌象特征类别",
    "  ▸ 训练集/验证集按比例划分",
]
add_multi(slide, cx + Inches(0.2), col_y + Inches(0.65), col_w - Inches(0.4), col_h - Inches(0.8),
          data_lines, fs=10, ls=1.5)

# 栏3: 数据标注
cx = sx + (col_w + col_gap) * 2
card(slide, cx, col_y, col_w, col_h)
add_rect(slide, cx, col_y, col_w, Inches(0.5), fill_color=VERMILION)
add_text(slide, cx + Inches(0.2), col_y + Inches(0.05), col_w, Inches(0.4),
         "03 数据标注", fs=14, color=WHITE, bold=True)
label_lines = [
    "标注工具：",
    "  ▸ LabelImg / Roboflow",
    "",
    "标注方式：",
    "  ▸ 边界框标注（Bounding Box）",
    "  ▸ 每张图片标注舌象特征区域",
    "",
    "标注类别（21类）：",
    "  ▸ 舌质：淡红舌、红舌、绛舌、",
    "    紫舌、淡白舌",
    "  ▸ 舌苔：薄白苔、厚白苔、",
    "    薄黄苔、厚黄苔、腻苔、",
    "    剥苔、无苔",
    "  ▸ 舌形：齿痕舌、胖大舌、",
    "    瘦薄舌、裂纹舌",
    "  ▸ 其他：舌下络脉、点刺舌等",
    "",
    "标注格式：YOLO txt格式",
    "  class x_center y_center w h",
]
add_multi(slide, cx + Inches(0.2), col_y + Inches(0.65), col_w - Inches(0.4), col_h - Inches(0.8),
          label_lines, fs=10, ls=1.5)

footer(slide, 6)


# Slide 7: 数据集制作 & 模型训练 & 模型验证
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "数据集制作 · 模型训练 · 模型验证", "构建数据集 → YOLO训练 → 效果验证", "02 产品设计过程")

col_w = Inches(3.9)
col_gap = Inches(0.2)
sx = MARGIN
col_y = Inches(1.5)
col_h = Inches(5.3)

# 栏1: 数据集制作
cx = sx
card(slide, cx, col_y, col_w, col_h)
add_rect(slide, cx, col_y, col_w, Inches(0.5), fill_color=JADE_DEEP)
add_text(slide, cx + Inches(0.2), col_y + Inches(0.05), col_w, Inches(0.4),
         "04 数据集制作", fs=14, color=WHITE, bold=True)
ds_lines = [
    "数据集划分：",
    "  ▸ 训练集（Train）约80%",
    "  ▸ 验证集（Valid）约15%",
    "  ▸ 测试集（Test）约5%",
    "",
    "数据增强：",
    "  ▸ 随机旋转（±15°）",
    "  ▸ 水平/垂直翻转",
    "  ▸ 色彩抖动（HSV）",
    "  ▸ 随机缩放裁剪",
    "  ▸ Mosaic拼接增强",
    "",
    "目录结构：",
    "  dataset/",
    "    ├── images/",
    "    │   ├── train/",
    "    │   └── val/",
    "    └── labels/",
    "        ├── train/",
    "        └── val/",
    "",
    "配置文件：data.yaml",
    "  nc: 21  # 类别数",
]
add_multi(slide, cx + Inches(0.2), col_y + Inches(0.65), col_w - Inches(0.4), col_h - Inches(0.8),
          ds_lines, fs=10, ls=1.5)

# 栏2: 模型训练
cx = sx + col_w + col_gap
card(slide, cx, col_y, col_w, col_h)
add_rect(slide, cx, col_y, col_w, Inches(0.5), fill_color=JADE)
add_text(slide, cx + Inches(0.2), col_y + Inches(0.05), col_w, Inches(0.4),
         "05 模型训练", fs=14, color=WHITE, bold=True)
train_lines = [
    "模型选择：",
    "  ▸ YOLOv8 — 目标检测",
    "  ▸ 单阶段检测，速度快",
    "  ▸ 适合实时舌象特征检测",
    "",
    "训练参数：",
    "  ▸ 预训练权重：yolov8s.pt",
    "  ▸ 训练轮数：100-300 epochs",
    "  ▸ 批次大小：16",
    "  ▸ 图像尺寸：640×640",
    "  ▸ 优化器：AdamW",
    "  ▸ 学习率：0.001",
    "",
    "训练过程：",
    "  ▸ 加载预训练模型",
    "  ▸ 冻结backbone微调",
    "  ▸ 渐进式解冻训练",
    "  ▸ Loss收敛监控",
    "",
    "输出：best.pt（最优权重）",
]
add_multi(slide, cx + Inches(0.2), col_y + Inches(0.65), col_w - Inches(0.4), col_h - Inches(0.8),
          train_lines, fs=10, ls=1.5)

# 栏3: 模型验证
cx = sx + (col_w + col_gap) * 2
card(slide, cx, col_y, col_w, col_h)
add_rect(slide, cx, col_y, col_w, Inches(0.5), fill_color=GOLD)
add_text(slide, cx + Inches(0.2), col_y + Inches(0.05), col_w, Inches(0.4),
         "06 模型验证", fs=14, color=WHITE, bold=True)
val_lines = [
    "评估指标：",
    "  ▸ mAP@0.5 — 平均精度均值",
    "  ▸ Precision — 精确率",
    "  ▸ Recall — 召回率",
    "  ▸ F1-Score — 综合指标",
    "",
    "验证结果：",
    "  ▸ 21类舌象特征均可检测",
    "  ▸ 主要特征检测准确率高",
    "  ▸ 部分相似特征需优化",
    "    （如薄白苔vs厚白苔）",
    "",
    "迭代优化：",
    "  ▸ 增加难例挖掘（Hard Mining）",
    "  ▸ 补充边缘样本数据",
    "  ▸ 调整置信度阈值",
    "  ▸ 多轮训练迭代",
    "",
    "最终模型：best.pt",
    "  部署至Flask后端本地调用",
]
add_multi(slide, cx + Inches(0.2), col_y + Inches(0.65), col_w - Inches(0.4), col_h - Inches(0.8),
          val_lines, fs=10, ls=1.5)

footer(slide, 7)


# Slide 8: 基础网页开发 & 智能体上线
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "基础网页开发 · 智能体功能上线", "Flask框架搭建 → 青囊Agent对话功能", "02 产品设计过程")

col_w = Inches(5.9)
col_gap = Inches(0.3)
sx = MARGIN
col_y = Inches(1.5)
col_h = Inches(5.3)

# 左：基础网页开发
cx = sx
card(slide, cx, col_y, col_w, col_h)
add_rect(slide, cx, col_y, col_w, Inches(0.5), fill_color=VERMILION)
add_text(slide, cx + Inches(0.2), col_y + Inches(0.05), col_w, Inches(0.4),
         "07 基础网页开发", fs=14, color=WHITE, bold=True)
web_lines = [
    "技术选型：",
    "  ▸ 后端：Python Flask 轻量框架",
    "  ▸ 前端：HTML + CSS + JavaScript",
    "  ▸ 数据库：SQLite（用户+记录）",
    "  ▸ 认证：Werkzeug PBKDF2 密码加密",
    "",
    "页面开发（8个核心页面）：",
    "  ▸ 首页 — 品牌展示 + 功能入口",
    "  ▸ 舌象分析 — 图片上传 + YOLO检测",
    "  ▸ 声音分析 — 录音 + 音频特征提取",
    "  ▸ 体质问卷 — 15项大白话量表",
    "  ▸ 面相分析 — 摄像头拍照 + 五色分析",
    "  ▸ AI报告 — DeepSeek综合分析",
    "  ▸ 用户仪表盘 — 历史记录管理",
    "  ▸ 管理后台 — 数据统计与管理",
    "",
    "UI设计：",
    "  ▸ 新中式水墨风：墨翠绿+宣纸白+赭石金",
    "  ▸ 望闻问切图片融入功能卡片",
    "  ▸ 入场动画视频 + 免责声明弹窗",
]
add_multi(slide, cx + Inches(0.2), col_y + Inches(0.65), col_w - Inches(0.4), col_h - Inches(0.8),
          web_lines, fs=10, ls=1.5)

# 右：智能体上线
cx = sx + col_w + col_gap
card(slide, cx, col_y, col_w, col_h)
add_rect(slide, cx, col_y, col_w, Inches(0.5), fill_color=JADE_DEEP)
add_text(slide, cx + Inches(0.2), col_y + Inches(0.05), col_w, Inches(0.4),
         "08 智能体功能上线", fs=14, color=WHITE, bold=True)
agent_lines = [
    "青囊Agent · 中医体质科普智能体：",
    "",
    "  身份定位：",
    "  ▸ 精通中医基础理论、九种体质、",
    "    八纲辨证、食疗药膳的科普智能体",
    "",
    "  核心能力：",
    "  ▸ 接收用户症状描述（睡眠、二便、",
    "    出汗、寒热、精神状态等）",
    "  ▸ 结合望闻问切检测数据分析体质",
    "  ▸ 输出三段式结构：体质解析 + ",
    "    日常调养 + 饮食建议",
    "",
    "  安全铁律：",
    "  ▸ 禁止确诊疾病、开处方、写剂量",
    "  ▸ 急症第一时间提醒就医",
    "  ▸ 固定追加免责声明",
    "",
    "  技术实现：",
    "  ▸ DeepSeek API 驱动",
    "  ▸ 系统提示词工程（Prompt Engineering）",
    "  ▸ 自动注入用户检测数据上下文",
    "  ▸ 历史对话记忆（最近10轮）",
]
add_multi(slide, cx + Inches(0.2), col_y + Inches(0.65), col_w - Inches(0.4), col_h - Inches(0.8),
          agent_lines, fs=10, ls=1.5)

footer(slide, 8)


# Slide 9: 商品广告功能上线 & 网站测试
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "商品广告功能上线 · 网站测试", "虚拟货架 + 食疗推荐 → 全面测试上线", "02 产品设计过程")

col_w = Inches(5.9)
col_gap = Inches(0.3)
sx = MARGIN
col_y = Inches(1.5)
col_h = Inches(5.3)

# 左：商品广告功能
cx = sx
card(slide, cx, col_y, col_w, col_h)
add_rect(slide, cx, col_y, col_w, Inches(0.5), fill_color=JADE)
add_text(slide, cx + Inches(0.2), col_y + Inches(0.05), col_w, Inches(0.4),
         "09 商品广告功能上线", fs=14, color=WHITE, bold=True)
ad_lines = [
    "食疗推荐模块：",
    "  ▸ AI报告生成后，自动调用DeepSeek",
    "    二次生成食疗推荐",
    "  ▸ 基于用户体质倾向精准推荐",
    "  ▸ 输出：适合多吃 / 建议少吃 /",
    "    推荐食疗方（2-3道家常菜）",
    "  ▸ 明确标注「非药品，日常食材参考」",
    "",
    "虚拟货架「青囊养生坊」：",
    "  ▸ 四季养生茶饮包 — 药食同源 · 体质定制",
    "  ▸ 体质食疗食材礼盒 — 针对体质精选搭配",
    "  ▸ 儿童健脾开胃糊 — 山药莲子芡实研磨",
    "  ▸ 节气养生汤包 — 二十四节气对应汤料",
    "",
    "商业化铺垫：",
    "  ▸ 所有商品标注「即将上线 · 敬请期待」",
    "  ▸ 根据体质分析精准推荐商品",
    "  ▸ 为未来电商接入预留接口",
    "  ▸ 首页增加望闻问切图片引导",
]
add_multi(slide, cx + Inches(0.2), col_y + Inches(0.65), col_w - Inches(0.4), col_h - Inches(0.8),
          ad_lines, fs=10, ls=1.5)

# 右：网站测试
cx = sx + col_w + col_gap
card(slide, cx, col_y, col_w, col_h)
add_rect(slide, cx, col_y, col_w, Inches(0.5), fill_color=GOLD)
add_text(slide, cx + Inches(0.2), col_y + Inches(0.05), col_w, Inches(0.4),
         "10 网站测试", fs=14, color=WHITE, bold=True)
test_lines = [
    "功能测试：",
    "  ▸ 用户注册/登录/登出流程",
    "  ▸ 舌象图片上传 + YOLO检测验证",
    "  ▸ 声音录制 + 音频特征分析",
    "  ▸ 体质问卷提交 + 评分计算",
    "  ▸ 面相拍照 + 五色分析",
    "  ▸ AI报告生成（三诊必做解锁）",
    "  ▸ 青囊Agent对话功能",
    "  ▸ 食疗推荐 + 虚拟货架展示",
    "  ▸ 管理员后台数据管理",
    "",
    "性能优化：",
    "  ▸ 图片压缩（JPEG质量75，最大1024px）",
    "  ▸ Session数据精简（避免Cookie超限）",
    "  ▸ 关闭debug模式（避免响应超限）",
    "  ▸ YOLO模型懒加载",
    "",
    "问题修复记录：",
    "  ▸ 中文引号导致Python语法错误",
    "  ▸ Session Cookie超4KB丢失数据",
    "  ▸ 端口占用 → 进程清理",
]
add_multi(slide, cx + Inches(0.2), col_y + Inches(0.65), col_w - Inches(0.4), col_h - Inches(0.8),
          test_lines, fs=10, ls=1.5)

footer(slide, 9)


# ============================================================
# ===== 第三章：产品展示 =====
# ============================================================

# Slide 10: 产品展示 — 功能介绍（望闻问切）
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "功能介绍：望闻问切四诊体系", "AI赋能传统中医四诊", "03 产品展示")

features = [
    ("望", "舌象分析", "YOLO模型检测21种舌象特征\n舌质/舌苔/舌形多维分析", "wang.jpg", JADE),
    ("闻", "声音分析", "录制语音AI分析声音特征\n力度/音调/清晰度/语速", "wen.jpg", GOLD),
    ("问", "体质问卷", "15项大白话体质量表\n九种体质分型判定", "wen_ask.jpg", VERMILION),
    ("切", "AI健康报告", "DeepSeek综合分析\n个性化健康建议+食疗推荐", "qie.jpg", JADE_DEEP),
]

cw = Inches(2.8)
ch = Inches(5.0)
gap = Inches(0.25)
total = cw * 4 + gap * 3
sx = (SLIDE_W - total) / 2
sy = Inches(1.5)

for i, (char, title, desc, img_name, color) in enumerate(features):
    cx = sx + i * (cw + gap)
    card(slide, cx, sy, cw, ch)
    add_rect(slide, cx, sy, cw, Inches(0.08), fill_color=color)
    img_path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, cx + Inches(0.3), sy + Inches(0.3),
                                 cw - Inches(0.6), Inches(2.3))
    add_text(slide, cx, sy + Inches(2.8), cw, Inches(0.6),
             char, fs=38, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, cx, sy + Inches(3.4), cw, Inches(0.4),
             title, fs=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, cx + Inches(0.2), sy + Inches(3.85), cw - Inches(0.4), Inches(0.9),
             desc, fs=10, color=MUTED, align=PP_ALIGN.CENTER, ls=1.5)

footer(slide, 10)


# Slide 11: 产品展示 — 青囊Agent & 商业模式
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "青囊Agent · 商业模式 · 亮点介绍", "智能体问诊 + 商业化闭环 + 产品亮点", "03 产品展示")

# 左：青囊Agent
lx = MARGIN
lw = Inches(5.8)
card(slide, lx, Inches(1.5), lw, Inches(2.5))
add_rect(slide, lx, Inches(1.5), lw, Inches(0.45), fill_color=JADE)
add_text(slide, lx + Inches(0.2), Inches(1.53), lw, Inches(0.4),
         "🩺 青囊Agent", fs=14, color=WHITE, bold=True)
agent_lines = [
    "  ▸ 对话式中医体质科普智能体",
    "  ▸ 自动注入望闻问切检测数据",
    "  ▸ 三段式输出：体质解析+调养+饮食",
    "  ▸ 快捷提问按钮，降低使用门槛",
    "  ▸ 历史对话记忆（最近10轮）",
]
add_multi(slide, lx + Inches(0.2), Inches(2.05), lw - Inches(0.4), Inches(1.8),
          agent_lines, fs=11, ls=1.6)

# 左下：商业模式
card(slide, lx, Inches(4.2), lw, Inches(2.6))
add_rect(slide, lx, Inches(4.2), lw, Inches(0.45), fill_color=GOLD)
add_text(slide, lx + Inches(0.2), Inches(4.23), lw, Inches(0.4),
         "💰 商业模式", fs=14, color=WHITE, bold=True)
biz_lines = [
    "当前阶段（免费引流）：",
    "  ▸ 免费体质检测 + AI报告",
    "  ▸ 免费青囊Agent对话咨询",
    "  ▸ 积累用户量和检测数据",
    "",
    "商业化路径（规划中）：",
    "  ▸ 虚拟货架商品销售（食疗食材礼盒等）",
    "  ▸ 会员订阅（深度报告+专属Agent）",
    "  ▸ 线下中医馆合作导流",
    "  ▸ 企业健康福利采购",
]
add_multi(slide, lx + Inches(0.2), Inches(4.75), lw - Inches(0.4), Inches(1.9),
          biz_lines, fs=10, ls=1.5)

# 右：亮点介绍
rx = Inches(7.1)
rw = Inches(5.8)
card(slide, rx, Inches(1.5), rw, Inches(5.3))
add_rect(slide, rx, Inches(1.5), rw, Inches(0.45), fill_color=VERMILION)
add_text(slide, rx + Inches(0.2), Inches(1.53), rw, Inches(0.4),
         "✨ 产品亮点", fs=14, color=WHITE, bold=True)

highlights = [
    ("🎯", "YOLO本地部署", "21种舌象特征实时检测，\n无需云端调用，响应快"),
    ("🤖", "DeepSeek AI驱动", "报告生成+食疗推荐+Agent\n对话，三大AI场景覆盖"),
    ("🎨", "新中式水墨UI", "墨翠绿+宣纸白+赭石金，\n望闻问切书法图片融入设计"),
    ("📖", "大白话问卷", "专业术语→通俗表达+对比参照，\n家长零门槛使用"),
    ("🛒", "商业化闭环", "检测→分析→建议→食疗→\n虚拟货架，完整商业链路"),
    ("🔐", "完整用户体系", "注册登录+数据隔离+\n管理员后台，全链路覆盖"),
]
for i, (icon, t, d) in enumerate(highlights):
    hy = Inches(2.1) + i * Inches(0.75)
    add_text(slide, rx + Inches(0.2), hy, Inches(0.4), Inches(0.5),
             icon, fs=18, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, rx + Inches(0.7), hy + Inches(0.02), Inches(2.2), Inches(0.3),
             t, fs=12, color=INK, bold=True)
    add_text(slide, rx + Inches(0.7), hy + Inches(0.32), Inches(4.5), Inches(0.4),
             d, fs=9, color=MUTED, ls=1.4)

footer(slide, 11)


# Slide 12: 产品展示 — 虚拟货架
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "虚拟货架：青囊养生坊", "基于体质分析的精准商品推荐", "03 产品展示")

products = [
    ("🍵", "四季养生茶饮包", "药食同源", "根据体质定制四季茶饮搭配方案", JADE),
    ("🥗", "体质食疗食材礼盒", "食材精选", "针对您体质倾向挑选的药食同源食材组合", GOLD),
    ("🥣", "儿童健脾开胃糊", "药膳辅食", "山药、莲子、芡实等健脾食材研磨", VERMILION),
    ("🍲", "节气养生汤包", "时令推荐", "二十四节气对应养生汤料搭配", JADE_DEEP),
]

pw = Inches(2.8)
gap = Inches(0.25)
total = pw * 4 + gap * 3
sx = (SLIDE_W - total) / 2
sy = Inches(1.6)

for i, (icon, name, tag, desc, color) in enumerate(products):
    px = sx + i * (pw + gap)
    card(slide, px, sy, pw, Inches(3.5))
    add_rect(slide, px, sy, pw, Inches(0.06), fill_color=color)
    add_text(slide, px, sy + Inches(0.25), pw, Inches(0.6),
             icon, fs=32, align=PP_ALIGN.CENTER)
    tag_pill(slide, px + Inches(0.75), sy + Inches(0.95), Inches(1.3), Inches(0.3), tag,
             bg=JADE_LIGHT, tc=color, fs=10)
    add_text(slide, px + Inches(0.15), sy + Inches(1.4), pw - Inches(0.3), Inches(0.5),
             name, fs=13, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, px + Inches(0.2), sy + Inches(1.9), pw - Inches(0.4), Inches(0.8),
             desc, fs=10, color=MUTED, align=PP_ALIGN.CENTER, ls=1.5)
    add_text(slide, px + Inches(0.15), sy + Inches(2.7), pw - Inches(0.3), Inches(0.3),
             "即将上架", fs=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, px + Inches(0.15), sy + Inches(3.05), pw - Inches(0.3), Inches(0.25),
             "即将上线 · 敬请期待", fs=9, color=MUTED, align=PP_ALIGN.CENTER)

# 底部说明
add_card = card(slide, MARGIN, Inches(5.5), Inches(12.1), Inches(1.2), fill=GOLD_LIGHT)
add_text(slide, MARGIN + Inches(0.3), Inches(5.6), Inches(11.5), Inches(0.4),
         "📌 商业化逻辑", fs=13, color=GOLD, bold=True)
add_multi(slide, MARGIN + Inches(0.3), Inches(5.95), Inches(11.5), Inches(0.7),
          ["用户完成体质检测 → AI分析体质倾向 → 精准推荐食疗食材 → 虚拟货架商品转化 → 未来电商闭环"],
          fs=12, color=INK)

footer(slide, 12)


# ============================================================
# ===== 第四章：未来展望 =====
# ============================================================

# Slide 13: 未来展望
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "未来展望", "功能规划 · 技术升级 · 商业化路径", "04 未来展望")

futures = [
    ("📱", "移动端适配", "开发小程序/APP版本\n支持微信扫码使用\n摄像头实时舌象拍摄", JADE),
    ("🧠", "模型升级", "扩充舌象数据集\n增加面诊AI模型\n提升检测准确率", GOLD),
    ("💊", "电商闭环", "虚拟商品正式上架\n接入支付系统\n体质定制食材配送", VERMILION),
    ("👨‍⚕️", "在线问诊", "对接真实中医师\n线上视频问诊\n电子处方流转", JADE_DEEP),
    ("📊", "健康档案", "长期体质追踪\n趋势变化图表\n家庭成员管理", JADE),
    ("🤝", "机构合作", "中医馆/体检中心\n企业健康福利\n社区健康服务站", GOLD),
]

cw = Inches(3.8)
ch = Inches(2.5)
gx = Inches(0.3)
gy = Inches(0.25)
cols = 3
sx = (SLIDE_W - (cw * cols + gx * (cols - 1))) / 2
sy = Inches(1.5)

for i, (icon, t, d, c) in enumerate(futures):
    col = i % cols
    row = i // cols
    cx = sx + col * (cw + gx)
    cy = sy + row * (ch + gy)
    card(slide, cx, cy, cw, ch)
    add_rect(slide, cx, cy, Inches(0.08), ch, fill_color=c)
    add_text(slide, cx + Inches(0.25), cy + Inches(0.2), Inches(0.6), Inches(0.5),
             icon, fs=24, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, cx + Inches(0.9), cy + Inches(0.2), Inches(2.5), Inches(0.4),
             t, fs=15, color=c, bold=True)
    add_text(slide, cx + Inches(0.9), cy + Inches(0.65), cw - Inches(1.1), Inches(1.5),
             d, fs=11, color=MUTED, ls=1.6)

footer(slide, 13)


# Slide 14: 未来展望 — 发展路线图
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_bar(slide, "发展路线图", "从MVP到平台的演进规划", "04 未来展望")

phases = [
    ("Phase 1", "已完成 ✅", "MVP最小可行产品", JADE,
     ["YOLO舌象检测", "四诊功能（望闻问切）", "AI健康报告", "青囊Agent对话", "虚拟货架展示", "用户系统+管理后台"]),
    ("Phase 2", "近期规划", "产品完善 & 商业化启动", GOLD,
     ["移动端适配（小程序）", "电商支付系统接入", "模型数据扩充升级", "会员订阅体系", "健康档案长期追踪"]),
    ("Phase 3", "中期目标", "平台化 & 生态拓展", VERMILION,
     ["在线中医师问诊", "线下中医馆合作", "企业健康福利采购", "社区健康服务站", "开放API平台"]),
]

pw = Inches(3.8)
gap = Inches(0.3)
total = pw * 3 + gap * 2
sx = (SLIDE_W - total) / 2
sy = Inches(1.5)
ph = Inches(5.0)

for i, (phase, status, title, color, items) in enumerate(phases):
    px = sx + i * (pw + gap)
    card(slide, px, sy, pw, ph)
    add_rect(slide, px, sy, pw, Inches(0.7), fill_color=color)
    add_text(slide, px + Inches(0.2), sy + Inches(0.05), pw, Inches(0.3),
             phase, fs=12, color=GOLD_LIGHT if color != GOLD else WHITE)
    add_text(slide, px + Inches(0.2), sy + Inches(0.35), pw, Inches(0.3),
             status, fs=11, color=WHITE, bold=True)
    add_text(slide, px + Inches(0.2), sy + Inches(0.85), pw, Inches(0.4),
             title, fs=14, color=INK, bold=True)

    item_lines = [f"  ▸  {item}" for item in items]
    add_multi(slide, px + Inches(0.15), sy + Inches(1.35), pw - Inches(0.3), Inches(3.3),
              item_lines, fs=11, color=MUTED, ls=1.7)

footer(slide, 14)


# ============================================================
# Slide 15: 结尾
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, JADE_DEEP)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=JADE_DEEP)

for cx, cy, r in [(Inches(11), Inches(1), Inches(2)), (Inches(2), Inches(6), Inches(1.5))]:
    add_rect(slide, cx - r, cy - r, r * 2, r * 2, fill_color=JADE, shape_type=MSO_SHAPE.OVAL)

add_rect(slide, Inches(4), Inches(2.3), Inches(5.333), Pt(2), fill_color=GOLD)

add_text(slide, Inches(2), Inches(2.6), Inches(9.333), Inches(1),
         "青囊AI", fs=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(3.7), Inches(9.333), Inches(0.5),
         "以AI之眼观舌象 · 以古人之智辨体质", fs=18, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(4.5), Inches(9.333), Inches(0.5),
         "感谢观看 · 欢迎体验", fs=16, color=GOLD, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(5.2), Inches(9.333), Inches(0.4),
         "访问地址：http://127.0.0.1:5000", fs=13,
         color=RGBColor(0xAA, 0xBB, 0xB5), align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(5.9), Inches(9.333), Inches(0.4),
         "⚠️ 本应用仅作中医体质健康科普参考，不构成医疗诊断", fs=11,
         color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# 保存
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "青囊AI项目演示.pptx")
prs.save(output_path)
print(f"PPT已保存至：{output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")

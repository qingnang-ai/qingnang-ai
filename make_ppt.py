# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static", "images")
OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "青囊AI项目大纲PPT.pptx")

JADE = RGBColor(0x0d, 0x6b, 0x5b)
JADE_DEEP = RGBColor(0x07, 0x3b, 0x35)
GOLD = RGBColor(0xc8, 0x9b, 0x47)
GOLD_LIGHT = RGBColor(0xf5, 0xed, 0xd6)
PAPER = RGBColor(0xfa, 0xf8, 0xf3)
INK = RGBColor(0x17, 0x33, 0x2f)
WHITE = RGBColor(0xff, 0xff, 0xff)
MUTED = RGBColor(0x6a, 0x7d, 0x78)
VERMILION = RGBColor(0xb5, 0x64, 0x58)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
SLIDE_W = 10
SLIDE_H = 5.625

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        srgb = shape.fill._fill.find(f'{{{ns}}}srgbClr')
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, f'{{{ns}}}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape

def add_text(slide, x, y, w, h, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multi_text(slide, x, y, w, h, lines, size=14, color=INK, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(size * 0.6)
    return txBox

def add_image_safe(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

def add_decor(slide):
    add_rect(slide, 0, 0, SLIDE_W, 0.08, JADE)
    add_rect(slide, 0, SLIDE_H - 0.08, SLIDE_W, 0.08, GOLD)

# ====== Slide 1: Cover ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, JADE_DEEP)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, JADE_DEEP)

add_rect(slide, 1.5, 1.8, 7, 0.04, GOLD)
add_text(slide, 1.5, 0.8, 7, 1, '青囊AI', size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name='Microsoft YaHei')
add_text(slide, 1.5, 2.1, 7, 0.5, '基于AI的中医体质智能分析平台', size=18, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
add_rect(slide, 1.5, 2.9, 7, 0.04, GOLD)
add_text(slide, 1.5, 3.2, 7, 0.4, '望 · 闻 · 问 · 切  |  AI × 中医', size=14, color=RGBColor(0xaa, 0xbb, 0xb5), align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.5, 7, 0.3, '2026', size=12, color=RGBColor(0x88, 0x99, 0x95), align=PP_ALIGN.CENTER)

# ====== Slide 2: TOC ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PAPER)
add_decor(slide)
add_text(slide, 0.8, 0.4, 4, 0.6, '目录', size=28, color=JADE_DEEP, bold=True)
add_rect(slide, 0.8, 1.05, 0.6, 0.04, GOLD)

toc_items = [
    ('01', '项目背景', '为什么做这件事'),
    ('02', '解决的问题', '痛点在哪里'),
    ('03', '我们的思路', '怎么做的'),
    ('04', '望 · 舌象分析', '拍照识别21种舌象特征'),
    ('05', '闻 · 声音分析', '录一段话分析体质'),
    ('06', '问 · 体质问卷', '大白话问卷'),
    ('07', '切 · 面相分析', '面部脏腑对应'),
    ('08', 'AI报告与Agent', '智能分析与对话'),
    ('09', '虚拟商城', '检测到消费闭环'),
]
for i, (num, title, desc) in enumerate(toc_items):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 3.0
    y = 1.4 + row * 1.3
    add_text(slide, x, y, 0.6, 0.5, num, size=22, color=GOLD, bold=True)
    add_text(slide, x + 0.6, y + 0.05, 2.2, 0.4, title, size=14, color=INK, bold=True)
    add_text(slide, x + 0.6, y + 0.45, 2.2, 0.35, desc, size=10, color=MUTED)

# ====== Slide 3: 项目背景 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PAPER)
add_decor(slide)
add_text(slide, 0.8, 0.4, 4, 0.6, '项目背景', size=28, color=JADE_DEEP, bold=True)
add_rect(slide, 0.8, 1.05, 0.6, 0.04, GOLD)

bg_box = add_rect(slide, 0.8, 1.4, 8.4, 3.5, WHITE)
bg_text = [
    '很多家长想了解孩子的体质，但去趟中医院并不方便，',
    '而且中医说的"阳虚""痰湿"这些词，普通人根本听不懂。',
    '',
    '市面上健康类App不少，但基本都是西医那套指标，',
    '真正把中医理论和AI结合起来、让老百姓用得动的产品，几乎没有。',
    '',
    '我们就想做这么一个东西。',
]
add_multi_text(slide, 1.2, 1.7, 7.6, 3, bg_text, size=15, color=INK, line_spacing=1.8)

# ====== Slide 4: 解决的问题 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PAPER)
add_decor(slide)
add_text(slide, 0.8, 0.4, 5, 0.6, '解决的问题', size=28, color=JADE_DEEP, bold=True)
add_rect(slide, 0.8, 1.05, 0.6, 0.04, GOLD)

problems = [
    ('不方便', '看中医不方便、太贵'),
    ('没工具', '舌象面色没有工具帮你看'),
    ('听不懂', '问卷全是专业术语，家长一头雾水'),
    ('没建议', '做完检测没人告诉你怎么吃、怎么养'),
]
for i, (tag, desc) in enumerate(problems):
    y = 1.4 + i * 0.9
    add_rect(slide, 0.8, y, 1.2, 0.6, JADE)
    add_text(slide, 0.8, y + 0.1, 1.2, 0.4, tag, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.1, y, 7.1, 0.6, RGBColor(0xf0, 0xf5, 0xf3))
    add_text(slide, 2.3, y + 0.12, 6.7, 0.4, desc, size=13, color=INK)

# ====== Slide 5: 我们的思路 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PAPER)
add_decor(slide)
add_text(slide, 0.8, 0.4, 5, 0.6, '我们的思路', size=28, color=JADE_DEEP, bold=True)
add_rect(slide, 0.8, 1.05, 0.6, 0.04, GOLD)

flow_items = ['望', '闻', '问', '切']
flow_colors = [JADE, RGBColor(0x2a, 0x8a, 0x7a), GOLD, VERMILION]
flow_descs = ['舌象分析', '声音分析', '体质问卷', '面相分析']
for i, (item, color, desc) in enumerate(zip(flow_items, flow_colors, flow_descs)):
    x = 0.8 + i * 2.2
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(1.6), Inches(1.2), Inches(1.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, x, 1.85, 1.2, 0.7, item, size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, 2.95, 1.2, 0.4, desc, size=12, color=INK, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow_x = x + 1.3
        add_text(slide, arrow_x, 1.9, 0.8, 0.5, '→', size=24, color=GOLD, align=PP_ALIGN.CENTER)

add_rect(slide, 0.8, 3.6, 8.4, 0.8, JADE_LIGHT := RGBColor(0xe7, 0xf4, 0xef))
add_text(slide, 1.0, 3.75, 8, 0.5, '把中医望闻问切搬上手机，用AI替代老中医的眼睛和经验', size=15, color=JADE_DEEP, bold=True, align=PP_ALIGN.CENTER)

# ====== Slide 6: 望 · 舌象分析 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PAPER)
add_decor(slide)
add_text(slide, 0.8, 0.4, 6, 0.6, '望 · 舌象分析', size=28, color=JADE_DEEP, bold=True)
add_rect(slide, 0.8, 1.05, 0.6, 0.04, GOLD)

add_image_safe(slide, os.path.join(IMG_DIR, "wang.jpg"), 0.8, 1.3, 3.5, 3.5)

add_text(slide, 4.6, 1.4, 5, 0.5, 'YOLO模型识别21种舌象特征', size=16, color=JADE_DEEP, bold=True)
features = [
    '拍张照就知道舌质颜色、舌苔厚度、舌形齿痕',
    '支持直接拍照或上传图片两种方式',
    '模型基于5270组标注数据训练',
    '自动标注舌象特征位置和类型',
    '检测结果作为体质分析的重要依据',
]
add_multi_text(slide, 4.6, 2.0, 5, 2.8, [f'• {f}' for f in features], size=12, color=INK, line_spacing=1.6)

# ====== Slide 7: 闻 · 声音分析 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PAPER)
add_decor(slide)
add_text(slide, 0.8, 0.4, 6, 0.6, '闻 · 声音分析', size=28, color=JADE_DEEP, bold=True)
add_rect(slide, 0.8, 1.05, 0.6, 0.04, GOLD)

add_image_safe(slide, os.path.join(IMG_DIR, "wen.jpg"), 0.8, 1.3, 3.5, 3.5)

add_text(slide, 4.6, 1.4, 5, 0.5, '录一段话，分析声音特征', size=16, color=JADE_DEEP, bold=True)
voice_features = [
    '音量大小（RMS能量）→ 判断中气足不足',
    '音调高低（基频）→ 判断体质寒热倾向',
    '声音清晰度（过零率）→ 辅助判断',
    '语速快慢（时长）→ 反映精神状态',
    '用户朗读指定文本，程序自动采集分析',
]
add_multi_text(slide, 4.6, 2.0, 5, 2.8, [f'• {f}' for f in voice_features], size=12, color=INK, line_spacing=1.6)

# ====== Slide 8: 问 · 体质问卷 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PAPER)
add_decor(slide)
add_text(slide, 0.8, 0.4, 6, 0.6, '问 · 体质问卷', size=28, color=JADE_DEEP, bold=True)
add_rect(slide, 0.8, 1.05, 0.6, 0.04, GOLD)

add_image_safe(slide, os.path.join(IMG_DIR, "wen_ask.jpg"), 0.8, 1.3, 3.5, 3.5)

add_text(slide, 4.6, 1.4, 5, 0.5, '全部改成大白话', size=16, color=JADE_DEEP, bold=True)
qa_examples = [
    '不说"山根发青"',
    '而是说"鼻梁中间皮肤发青或能看到青色血管，',
    '  跟两颊比偏青偏暗"',
    '',
    '不说"面色无华"',
    '而是说"脸色比起手臂内侧皮肤偏暗偏黄"',
    '',
    '每道题都有具体对比参照，家长一看就懂',
]
add_multi_text(slide, 4.6, 2.0, 5, 2.8, qa_examples, size=12, color=INK, line_spacing=1.5)

# ====== Slide 9: 切 · 面相分析 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PAPER)
add_decor(slide)
add_text(slide, 0.8, 0.4, 6, 0.6, '切 · 面相分析', size=28, color=JADE_DEEP, bold=True)
add_rect(slide, 0.8, 1.05, 0.6, 0.04, GOLD)

add_image_safe(slide, os.path.join(IMG_DIR, "qie.jpg"), 0.8, 1.3, 3.5, 3.5)

add_text(slide, 4.6, 1.4, 5, 0.5, '基于面部脏腑对应理论', size=16, color=JADE_DEEP, bold=True)
face_features = [
    '额头 → 候心肺',
    '鼻部 → 候脾胃',
    '左颊 → 候肝',
    '右颊 → 候肺',
    '下颌 → 候肾',
    '',
    '五色（青赤黄白黑）映射体质偏颇',
    '拍照即可分析，自动识别面部区域',
]
add_multi_text(slide, 4.6, 2.0, 5, 2.8, [f'• {f}' if f else f for f in face_features], size=12, color=INK, line_spacing=1.5)

# ====== Slide 10: AI报告与Agent ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PAPER)
add_decor(slide)
add_text(slide, 0.8, 0.4, 8, 0.6, 'AI综合报告 + 青囊Agent', size=28, color=JADE_DEEP, bold=True)
add_rect(slide, 0.8, 1.05, 0.6, 0.04, GOLD)

# Left: AI报告
add_rect(slide, 0.8, 1.4, 4.1, 3.5, WHITE)
add_rect(slide, 0.8, 1.4, 4.1, 0.5, JADE)
add_text(slide, 0.8, 1.5, 4.1, 0.3, 'AI综合报告', size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
report_items = [
    '整合望闻问切四诊数据',
    '调用DeepSeek大模型',
    '生成个性化体质分析报告',
    '包含食疗药膳建议',
    '基于体质推荐食材',
    '明确标注非药品、仅供参考',
]
add_multi_text(slide, 1.1, 2.1, 3.5, 2.5, [f'• {f}' for f in report_items], size=11, color=INK, line_spacing=1.5)

# Right: Agent
add_rect(slide, 5.1, 1.4, 4.1, 3.5, WHITE)
add_rect(slide, 5.1, 1.4, 4.1, 0.5, GOLD)
add_text(slide, 5.1, 1.5, 4.1, 0.3, '青囊Agent', size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
agent_items = [
    '像聊天一样问它问题',
    '基于九种体质、八纲辨证理论',
    '结合用户已有检测数据分析',
    '给出日常调养和饮食建议',
    '不确诊疾病、不开处方',
    '提醒急症及时就医',
]
add_multi_text(slide, 5.4, 2.1, 3.5, 2.5, [f'• {f}' for f in agent_items], size=11, color=INK, line_spacing=1.5)

# ====== Slide 11: 虚拟商城 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PAPER)
add_decor(slide)
add_text(slide, 0.8, 0.4, 6, 0.6, '虚拟商城', size=28, color=JADE_DEEP, bold=True)
add_rect(slide, 0.8, 1.05, 0.6, 0.04, GOLD)

add_text(slide, 0.8, 1.3, 8, 0.5, '根据体质推荐食材，从检测到调养到消费的闭环', size=15, color=JADE_DEEP, bold=True)

products = [
    ('四季养生茶饮包', '药食同源食材'),
    ('体质食疗食材礼盒', '按体质定制'),
    ('儿童健脾开胃糊', '适合儿童体质'),
    ('节气养生汤包', '应季推荐'),
]
for i, (name, tag) in enumerate(products):
    x = 0.8 + i * 2.2
    add_rect(slide, x, 2.0, 2.0, 2.2, WHITE)
    add_rect(slide, x, 2.0, 2.0, 0.6, GOLD_LIGHT)
    add_text(slide, x, 2.1, 2.0, 0.4, name, size=11, color=JADE_DEEP, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, 2.8, 2.0, 0.3, tag, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 3.5, 2.0, 0.5, RGBColor(0xe7, 0xf4, 0xef))
    add_text(slide, x, 3.6, 2.0, 0.3, '即将上线', size=10, color=JADE, align=PP_ALIGN.CENTER)

add_text(slide, 0.8, 4.5, 8.4, 0.4, '检测 → 分析 → 调养建议 → 推荐食材 → 消费闭环', size=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# ====== Slide 12: 结尾 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, JADE_DEEP)
add_text(slide, 1.5, 1.8, 7, 1, '谢谢', size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(slide, 3.5, 3.0, 3, 0.04, GOLD)
add_text(slide, 1.5, 3.3, 7, 0.4, '青囊AI · 望闻问切 · AI × 中医', size=14, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f"PPT saved: {OUTPUT}")

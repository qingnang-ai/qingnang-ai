# -*- coding: utf-8 -*-
"""基于路演模板生成青囊AI演示PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

TEMPLATE = r'c:\Users\66496\.trae-cn\attachments\6a7ec78c2d0cfdb32a600535\a4960b04-6656-4941-8dad-52b81a3908f8_8743cdee-25ad-4054-ab14-13371a8358f6_营员路演PPT.pptx'
IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static", "images")
OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "青囊AI路演PPT.pptx")

# 模板配色
PURPLE = RGBColor(0x4F, 0x37, 0xCD)
PURPLE_DARK = RGBColor(0x4A, 0x31, 0xCB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

prs = Presentation(TEMPLATE)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# 获取布局
layout_title = prs.slide_layouts[0]    # 标题和内容（封面）
layout_content = prs.slide_layouts[1]  # 标题和内容
layout_section = prs.slide_layouts[2]  # 节标题
# 如果有第4个布局
if len(prs.slide_layouts) > 3:
    layout_close = prs.slide_layouts[3]
else:
    layout_close = prs.slide_layouts[1]

# ============================================================
# 工具函数
# ============================================================
def set_text(shape, text, font_size=None, color=None, bold=None, font_name=None, align=None):
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.word_wrap = True
        for i, para in enumerate(tf.paragraphs):
            if i == 0:
                p = para
            else:
                p = tf.add_paragraph()
            p.text = text
            if font_size:
                p.font.size = Pt(font_size)
            if color:
                p.font.color.rgb = color
            if bold is not None:
                p.font.bold = bold
            if font_name:
                p.font.name = font_name
            if align:
                p.alignment = align
            break

def add_textbox(slide, x, y, w, h, text, fs=14, color=BLACK, bold=False, fn="微软雅黑",
                align=PP_ALIGN.LEFT, ls=1.5):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fs)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = fn
    p.alignment = align
    p.line_spacing = ls
    return tb

def add_multiline(slide, x, y, w, h, lines, fs=13, color=BLACK, fn="微软雅黑",
                  align=PP_ALIGN.LEFT, ls=1.6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fs)
        p.font.color.rgb = color
        p.font.name = fn
        p.alignment = align
        p.line_spacing = ls
        p.space_after = Pt(4)
    return tb

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_card(slide, x, y, w, h, fill=WHITE):
    s = add_rect(slide, x, y, w, h, fill_color=fill)
    s.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    s.line.width = Pt(0.75)
    return s

def add_title_accent(slide, title_text, subtitle=""):
    """添加带紫色竖条的标题"""
    add_rect(slide, Cm(1.2), Cm(0.6), Cm(0.25), Cm(1.2), fill_color=PURPLE_DARK)
    add_textbox(slide, Cm(1.7), Cm(0.55), Cm(28), Cm(1.0),
                title_text, fs=24, color=PURPLE, bold=True, fn="微软雅黑")
    if subtitle:
        add_textbox(slide, Cm(1.7), Cm(1.5), Cm(28), Cm(0.6),
                    subtitle, fs=12, color=GRAY, fn="微软雅黑")

def add_footer(slide, n):
    add_textbox(slide, Cm(1.2), Cm(17.5), Cm(15), Cm(0.5),
                "青囊AI · 中医体质智能分析平台", fs=8, color=GRAY)
    add_textbox(slide, Cm(33), Cm(17.5), Cm(1.5), Cm(0.5),
                str(n), fs=8, color=GRAY, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 1: 封面（修改模板第1页）
# ============================================================
slide = prs.slides[0]
for shape in slide.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if "项目名称" in run.text:
                    run.text = "青囊AI"
                elif "项目方向" in run.text:
                    run.text = "项目方向：AI + 中医体质分析"
                elif "汇报人" in run.text:
                    run.text = "汇报人：青囊AI团队"

# 在封面添加副标题
add_textbox(slide, Cm(8), Cm(11), Cm(20), Cm(1),
            "取意「青囊书」——华佗所传中医经典之名", fs=12, color=GRAY,
            align=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: 目录（新增）
# ============================================================
slide = prs.slides.add_slide(layout_content)
add_title_accent(slide, "目录 CONTENTS")

toc = [
    ("01", "项目背景", "市场需求 · 痛点分析 · 项目定位"),
    ("02", "产品设计过程", "需求分析→数据制作→模型训练→开发上线"),
    ("03", "产品展示", "功能介绍 · 商业模式 · 亮点介绍"),
    ("04", "未来展望", "功能规划 · 发展路线图"),
]

cw = Cm(15)
ch = Cm(3)
gap = Cm(0.5)
sx = Cm(1.7)
sy = Cm(3.5)

for i, (num, t, d) in enumerate(toc):
    cy = sy + i * (ch + gap)
    add_card(slide, sx, cy, cw, ch, fill=LIGHT_GRAY)
    add_rect(slide, sx, cy, Cm(0.2), ch, fill_color=PURPLE)
    add_textbox(slide, sx + Cm(0.5), cy + Cm(0.3), Cm(2), Cm(1),
                num, fs=28, color=PURPLE, bold=True, fn="Arial")
    add_textbox(slide, sx + Cm(3), cy + Cm(0.4), Cm(5), Cm(0.7),
                t, fs=16, color=BLACK, bold=True)
    add_textbox(slide, sx + Cm(3), cy + Cm(1.3), Cm(11), Cm(0.7),
                d, fs=11, color=GRAY)

add_footer(slide, 2)


# ============================================================
# Slide 3: 节标题 — 项目背景
# ============================================================
slide = prs.slides.add_slide(layout_section)
for shape in slide.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "单击此处" in run.text:
                    run.text = "01  项目背景"

add_textbox(slide, Cm(1.7), Cm(4), Cm(25), Cm(1),
            "市场需求 · 痛点分析 · 项目定位", fs=14, color=GRAY)


# ============================================================
# Slide 4: 市场需求与痛点
# ============================================================
slide = prs.slides.add_slide(layout_content)
add_title_accent(slide, "市场需求与痛点分析", "为什么需要青囊AI？")

# 左：痛点
add_textbox(slide, Cm(1.7), Cm(3), Cm(15), Cm(0.7),
            "当前痛点", fs=16, color=RGBColor(0xCC, 0x44, 0x44), bold=True)
pains = [
    "▸  中医门槛高：普通人难以理解体质概念",
    "    和舌象、面相等诊断指标的含义",
    "",
    "▸  问诊不便：线下中医挂号难、排队久，",
    "    儿童体质调理需求难以满足",
    "",
    "▸  缺乏工具：市面上缺少结合AI技术的",
    "    中医体质科普工具，家长无从下手",
    "",
    "▸  信息碎片化：网络健康信息良莠不齐，",
    "    缺乏系统性体质分析和调理建议",
]
add_multiline(slide, Cm(1.7), Cm(3.8), Cm(15), Cm(6), pains, fs=11, ls=1.5)

# 右：需求
rx = Cm(18)
add_textbox(slide, rx, Cm(3), Cm(15), Cm(0.7),
            "用户需求", fs=16, color=PURPLE, bold=True)
needs = [
    ("便捷检测", "在家即可完成舌象/声音/面相/问卷检测"),
    ("AI分析", "AI自动分析体质倾向，生成个性化建议"),
    ("通俗解读", "用大白话解释体质，家长看得懂"),
    ("调理方案", "提供食疗推荐和养生产品，形成闭环"),
]
for i, (t, d) in enumerate(needs):
    ny = Cm(3.8) + i * Cm(2)
    add_card(slide, rx, ny, Cm(15), Cm(1.7), fill=RGBColor(0xF0, 0xED, 0xFA))
    add_textbox(slide, rx + Cm(0.5), ny + Cm(0.15), Cm(5), Cm(0.5),
                t, fs=13, color=PURPLE, bold=True)
    add_textbox(slide, rx + Cm(0.5), ny + Cm(0.8), Cm(14), Cm(0.7),
                d, fs=10, color=GRAY)

add_footer(slide, 4)


# ============================================================
# Slide 5: 项目定位
# ============================================================
slide = prs.slides.add_slide(layout_content)
add_title_accent(slide, "项目定位与理念", "青囊AI是什么？")

# 中心定位卡片
add_card(slide, Cm(1.7), Cm(2.8), Cm(30.5), Cm(2), fill=PURPLE_DARK)
add_textbox(slide, Cm(2), Cm(2.9), Cm(30), Cm(0.6),
            "青囊AI · 中医体质智能分析平台", fs=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Cm(2), Cm(3.6), Cm(30), Cm(1),
            "结合传统中医望闻问切四诊理论，通过YOLO目标检测、DeepSeek AI大模型等技术，\n为普通家庭提供便捷的中医体质分析和调养建议，让千年医道智慧走进每个家庭。",
            fs=11, color=RGBColor(0xDD, 0xCC, 0xFF), align=PP_ALIGN.CENTER, ls=1.6)

# 四大支柱
pillars = [
    ("传承", "取意华佗《青囊书》\n传承千年中医智慧"),
    ("科技", "YOLO + DeepSeek\nAI赋能体质分析"),
    ("普惠", "面向普通家庭\n大白话解读体质"),
    ("闭环", "检测→分析→建议→\n食疗推荐→产品"),
]
pw = Cm(7)
gap = Cm(0.5)
total = pw * 4 + gap * 3
sx = (SLIDE_W - total) / 2
for i, (t, d) in enumerate(pillars):
    px = sx + i * (pw + gap)
    py = Cm(5.3)
    add_card(slide, px, py, pw, Cm(3))
    add_rect(slide, px, py, pw, Cm(0.15), fill_color=PURPLE)
    add_textbox(slide, px, py + Cm(0.4), pw, Cm(0.6),
                t, fs=16, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, px + Cm(0.3), py + Cm(1.1), pw - Cm(0.6), Cm(1.5),
                d, fs=10, color=GRAY, align=PP_ALIGN.CENTER, ls=1.5)

add_footer(slide, 5)


# ============================================================
# Slide 6: 节标题 — 产品设计过程
# ============================================================
slide = prs.slides.add_slide(layout_section)
for shape in slide.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "单击此处" in run.text:
                    run.text = "02  产品设计过程"

add_textbox(slide, Cm(1.7), Cm(4), Cm(25), Cm(1),
            "需求分析 → 数据制作 → 模型训练 → 开发上线", fs=14, color=GRAY)


# ============================================================
# Slide 7: 设计过程总览
# ============================================================
slide = prs.slides.add_slide(layout_content)
add_title_accent(slide, "产品设计过程 · 总览", "从需求到上线的完整研发流程")

steps = [
    ("01", "需求分析", "用户调研\n功能规划"),
    ("02", "数据收集", "舌象图片\n采集整理"),
    ("03", "数据标注", "21类特征\n标注标注"),
    ("04", "数据集制作", "训练/验证集\n数据增强"),
    ("05", "模型训练", "YOLOv8训练\n超参调优"),
    ("06", "模型验证", "mAP评估\n迭代优化"),
    ("07", "基础网页开发", "Flask框架\n中国风UI"),
    ("08", "智能体上线", "青囊Agent\nDeepSeek"),
    ("09", "商品广告上线", "虚拟货架\n食疗推荐"),
    ("10", "网站测试", "功能测试\n性能优化"),
]

cw = Cm(6)
ch = Cm(2.2)
gx = Cm(0.3)
gy = Cm(0.3)
cols = 5
total_w = cw * cols + gx * (cols - 1)
sx = (SLIDE_W - total_w) / 2
sy = Cm(2.8)

for i, (num, t, d) in enumerate(steps):
    col = i % cols
    row = i // cols
    cx = sx + col * (cw + gx)
    cy = sy + row * (ch + gy)
    add_card(slide, cx, cy, cw, ch, fill=WHITE)
    add_rect(slide, cx, cy, cw, Cm(0.12), fill_color=PURPLE)
    add_textbox(slide, cx + Cm(0.3), cy + Cm(0.2), Cm(2), Cm(0.5),
                num, fs=13, color=PURPLE, bold=True, fn="Arial")
    add_textbox(slide, cx + Cm(0.3), cy + Cm(0.8), cw - Cm(0.6), Cm(0.5),
                t, fs=12, color=BLACK, bold=True)
    add_textbox(slide, cx + Cm(0.3), cy + Cm(1.3), cw - Cm(0.6), Cm(0.8),
                d, fs=9, color=GRAY, ls=1.4)

add_textbox(slide, Cm(2), Cm(8.5), Cm(30), Cm(0.5),
            "完整研发周期：需求分析 → 数据工程 → 模型训练 → Web开发 → AI功能 → 商业化 → 测试上线",
            fs=11, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

add_footer(slide, 7)


# ============================================================
# Slide 8: 需求分析 + 数据收集 + 数据标注
# ============================================================
slide = prs.slides.add_slide(layout_content)
add_title_accent(slide, "需求分析 · 数据收集 · 数据标注", "明确目标 → 采集数据 → 标注特征")

col_w = Cm(10.5)
gap = Cm(0.4)
sx = Cm(1.5)
sy = Cm(2.8)
ch = Cm(13.5)

cols_data = [
    ("01 需求分析", PURPLE, [
        "目标用户：",
        "  ▸ 关注儿童体质的家长",
        "  ▸ 对中医养生感兴趣的成年人",
        "",
        "核心功能：",
        "  ▸ 舌象AI检测（望）",
        "  ▸ 声音特征分析（闻）",
        "  ▸ 体质问卷分型（问）",
        "  ▸ 面相分析（切）",
        "  ▸ AI综合健康报告",
        "  ▸ 中医体质科普Agent",
        "",
        "设计原则：",
        "  ▸ 中国风水墨UI",
        "  ▸ 大白话替代术语",
        "  ▸ 页面跳转简洁操作",
    ]),
    ("02 数据收集", RGBColor(0xCC, 0x88, 0x00), [
        "舌象图片采集：",
        "  ▸ 不同舌质颜色",
        "    （淡红、红、绛、紫、淡白）",
        "  ▸ 不同舌苔特征",
        "    （薄白、厚白、薄黄、厚黄、",
        "      腻苔、剥苔、无苔）",
        "  ▸ 舌形特征",
        "    （齿痕、胖大、瘦薄、裂纹）",
        "",
        "数据来源：",
        "  ▸ 公开舌象数据集",
        "  ▸ 中医舌诊图谱参考",
        "  ▸ 网络医学教育资源",
        "",
        "覆盖21种舌象特征类别",
    ]),
    ("03 数据标注", RGBColor(0xCC, 0x44, 0x44), [
        "标注工具：",
        "  ▸ LabelImg / Roboflow",
        "",
        "标注方式：",
        "  ▸ 边界框标注（Bounding Box）",
        "  ▸ 每张图片标注舌象特征区域",
        "",
        "标注类别（21类）：",
        "  ▸ 舌质5类：淡红/红/绛/紫/淡白",
        "  ▸ 舌苔7类：薄白/厚白/薄黄/厚黄",
        "    /腻苔/剥苔/无苔",
        "  ▸ 舌形4类：齿痕/胖大/瘦薄/裂纹",
        "  ▸ 其他5类：舌下络脉/点刺等",
        "",
        "标注格式：YOLO txt格式",
        "  class x y w h",
    ]),
]

for i, (title, color, lines) in enumerate(cols_data):
    cx = sx + i * (col_w + gap)
    add_card(slide, cx, sy, col_w, ch, fill=WHITE)
    add_rect(slide, cx, sy, col_w, Cm(0.6), fill_color=color)
    add_textbox(slide, cx + Cm(0.3), sy + Cm(0.05), col_w, Cm(0.5),
                title, fs=13, color=WHITE, bold=True)
    add_multiline(slide, cx + Cm(0.3), sy + Cm(0.8), col_w - Cm(0.6), ch - Cm(1),
                  lines, fs=9, ls=1.45)

add_footer(slide, 8)


# ============================================================
# Slide 9: 模型训练 + 验证 + 网页开发
# ============================================================
slide = prs.slides.add_slide(layout_content)
add_title_accent(slide, "模型训练 · 验证 · 网页开发 · 上线", "构建模型 → 验证优化 → 开发部署")

col_w = Cm(10.5)
gap = Cm(0.4)
sx = Cm(1.5)
sy = Cm(2.8)
ch = Cm(13.5)

cols_data = [
    ("04-06 模型训练与验证", PURPLE, [
        "模型选择：",
        "  ▸ YOLOv8 — 目标检测",
        "  ▸ 单阶段检测，速度快",
        "",
        "训练参数：",
        "  ▸ 预训练权重：yolov8s.pt",
        "  ▸ 训练轮数：100-300 epochs",
        "  ▸ 批次大小：16",
        "  ▸ 图像尺寸：640x640",
        "  ▸ 优化器：AdamW",
        "",
        "评估指标：",
        "  ▸ mAP@0.5 — 平均精度均值",
        "  ▸ Precision — 精确率",
        "  ▸ Recall — 召回率",
        "",
        "迭代优化：",
        "  ▸ 增加难例挖掘",
        "  ▸ 补充边缘样本数据",
        "  ▸ 最终输出：best.pt",
    ]),
    ("07 基础网页开发", RGBColor(0xCC, 0x88, 0x00), [
        "技术选型：",
        "  ▸ 后端：Python Flask",
        "  ▸ 前端：HTML+CSS+JS",
        "  ▸ 数据库：SQLite",
        "  ▸ 认证：Werkzeug加密",
        "",
        "核心页面（8个）：",
        "  ▸ 首页 — 品牌展示+功能入口",
        "  ▸ 舌象分析 — 上传+YOLO检测",
        "  ▸ 声音分析 — 录音+特征提取",
        "  ▸ 体质问卷 — 15项大白话量表",
        "  ▸ 面相分析 — 拍照+五色分析",
        "  ▸ AI报告 — DeepSeek分析",
        "  ▸ 用户仪表盘 — 记录管理",
        "  ▸ 管理后台 — 数据统计",
        "",
        "UI设计：新中式水墨风",
        "  墨翠绿+宣纸白+赭石金",
    ]),
    ("08-10 AI功能上线与测试", RGBColor(0xCC, 0x44, 0x44), [
        "青囊Agent上线：",
        "  ▸ DeepSeek API驱动",
        "  ▸ 系统提示词工程",
        "  ▸ 自动注入检测数据",
        "  ▸ 三段式输出+安全铁律",
        "",
        "商品广告功能：",
        "  ▸ 食疗推荐（AI二次生成）",
        "  ▸ 虚拟货架（4个商品）",
        "  ▸ 支付宝支付集成",
        "  ▸ 订单管理系统",
        "",
        "网站测试：",
        "  ▸ 功能测试（全流程）",
        "  ▸ 性能优化（图片压缩）",
        "  ▸ Bug修复（Cookie超限等）",
        "  ▸ 问卷优化（术语通俗化）",
        "",
        "部署：http://127.0.0.1:5000",
    ]),
]

for i, (title, color, lines) in enumerate(cols_data):
    cx = sx + i * (col_w + gap)
    add_card(slide, cx, sy, col_w, ch, fill=WHITE)
    add_rect(slide, cx, sy, col_w, Cm(0.6), fill_color=color)
    add_textbox(slide, cx + Cm(0.3), sy + Cm(0.05), col_w, Cm(0.5),
                title, fs=12, color=WHITE, bold=True)
    add_multiline(slide, cx + Cm(0.3), sy + Cm(0.8), col_w - Cm(0.6), ch - Cm(1),
                  lines, fs=9, ls=1.45)

add_footer(slide, 9)


# ============================================================
# Slide 10: 节标题 — 产品展示
# ============================================================
slide = prs.slides.add_slide(layout_section)
for shape in slide.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "单击此处" in run.text:
                    run.text = "03  产品展示"

add_textbox(slide, Cm(1.7), Cm(4), Cm(25), Cm(1),
            "功能介绍 · 商业模式 · 亮点介绍", fs=14, color=GRAY)


# ============================================================
# Slide 11: 功能介绍 — 望闻问切
# ============================================================
slide = prs.slides.add_slide(layout_content)
add_title_accent(slide, "功能介绍：望闻问切四诊体系", "AI赋能传统中医四诊")

features = [
    ("望", "舌象分析", "YOLO模型检测21种舌象特征\n舌质/舌苔/舌形多维分析", "wang.jpg"),
    ("闻", "声音分析", "录制语音AI分析声音特征\n力度/音调/清晰度/语速", "wen.jpg"),
    ("问", "体质问卷", "15项大白话体质量表\n九种体质分型判定", "wen_ask.jpg"),
    ("切", "AI健康报告", "DeepSeek综合分析\n个性化健康建议+食疗推荐", "qie.jpg"),
]

cw = Cm(7.5)
ch = Cm(12)
gap = Cm(0.5)
total = cw * 4 + gap * 3
sx = (SLIDE_W - total) / 2
sy = Cm(2.5)

for i, (char, title, desc, img_name) in enumerate(features):
    cx = sx + i * (cw + gap)
    add_card(slide, cx, sy, cw, ch, fill=WHITE)
    add_rect(slide, cx, sy, cw, Cm(0.15), fill_color=PURPLE)
    img_path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, cx + Cm(0.5), sy + Cm(0.5),
                                 cw - Cm(1), Cm(5.5))
    add_textbox(slide, cx, sy + Cm(6.3), cw, Cm(1),
                char, fs=32, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, cx, sy + Cm(7.8), cw, Cm(0.5),
                title, fs=14, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, cx + Cm(0.3), sy + Cm(8.5), cw - Cm(0.6), Cm(1.5),
                desc, fs=10, color=GRAY, align=PP_ALIGN.CENTER, ls=1.5)

add_footer(slide, 11)


# ============================================================
# Slide 12: 商业模式 & 亮点
# ============================================================
slide = prs.slides.add_slide(layout_content)
add_title_accent(slide, "商业模式 · 产品亮点", "商业化闭环 + 核心亮点")

# 左：商业模式
lx = Cm(1.5)
lw = Cm(15.5)
add_textbox(slide, lx, Cm(2.5), lw, Cm(0.6),
            "商业模式", fs=16, color=PURPLE, bold=True)

biz_lines = [
    "当前阶段（免费引流）：",
    "  ▸ 免费体质检测 + AI报告",
    "  ▸ 免费青囊Agent对话咨询",
    "  ▸ 积累用户量和检测数据",
    "",
    "商业化路径（规划中）：",
    "  ▸ 虚拟货架商品销售",
    "    （食疗食材礼盒等）",
    "  ▸ 支付宝支付已集成",
    "  ▸ 会员订阅（深度报告）",
    "  ▸ 线下中医馆合作导流",
    "  ▸ 企业健康福利采购",
]
add_multiline(slide, lx, Cm(3.2), lw, Cm(7), biz_lines, fs=10, ls=1.5)

# 右：亮点
rx = Cm(18)
rw = Cm(15.5)
add_textbox(slide, rx, Cm(2.5), rw, Cm(0.6),
            "产品亮点", fs=16, color=PURPLE, bold=True)

highlights = [
    ("YOLO本地部署", "21种舌象特征实时检测，无需云端"),
    ("DeepSeek AI驱动", "报告+食疗+Agent，三大AI场景"),
    ("新中式水墨UI", "墨翠绿+宣纸白，望闻问切图片融入"),
    ("大白话问卷", "术语通俗化+对比参照，零门槛"),
    ("商业化闭环", "检测→分析→建议→食疗→货架"),
    ("完整用户体系", "注册登录+数据隔离+管理后台"),
]
for i, (t, d) in enumerate(highlights):
    hy = Cm(3.2) + i * Cm(1.1)
    add_card(slide, rx, hy, rw, Cm(1), fill=RGBColor(0xF5, 0xF2, 0xFF))
    add_rect(slide, rx, hy, Cm(0.15), Cm(1), fill_color=PURPLE)
    add_textbox(slide, rx + Cm(0.4), hy + Cm(0.05), Cm(6), Cm(0.4),
                t, fs=11, color=BLACK, bold=True)
    add_textbox(slide, rx + Cm(0.4), hy + Cm(0.5), Cm(14), Cm(0.4),
                d, fs=9, color=GRAY)

add_footer(slide, 12)


# ============================================================
# Slide 13: 节标题 — 未来展望
# ============================================================
slide = prs.slides.add_slide(layout_section)
for shape in slide.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "单击此处" in run.text:
                    run.text = "04  未来展望"

add_textbox(slide, Cm(1.7), Cm(4), Cm(25), Cm(1),
            "功能规划 · 发展路线图", fs=14, color=GRAY)


# ============================================================
# Slide 14: 未来展望
# ============================================================
slide = prs.slides.add_slide(layout_content)
add_title_accent(slide, "未来展望", "功能规划 · 技术升级 · 商业化路径")

futures = [
    ("移动端适配", "开发小程序/APP版本\n支持微信扫码使用\n摄像头实时拍摄"),
    ("模型升级", "扩充舌象数据集\n增加面诊AI模型\n提升检测准确率"),
    ("电商闭环", "商品正式上架\n接入支付系统\n体质定制配送"),
    ("在线问诊", "对接真实中医师\n线上视频问诊\n电子处方流转"),
    ("健康档案", "长期体质追踪\n趋势变化图表\n家庭成员管理"),
    ("机构合作", "中医馆/体检中心\n企业健康福利\n社区健康服务"),
]

cw = Cm(10)
ch = Cm(3)
gx = Cm(0.5)
gy = Cm(0.4)
cols = 3
total_w = cw * cols + gx * (cols - 1)
sx = (SLIDE_W - total_w) / 2
sy = Cm(2.8)

for i, (t, d) in enumerate(futures):
    col = i % cols
    row = i // cols
    cx = sx + col * (cw + gx)
    cy = sy + row * (ch + gy)
    add_card(slide, cx, cy, cw, ch, fill=WHITE)
    add_rect(slide, cx, cy, Cm(0.2), ch, fill_color=PURPLE)
    add_textbox(slide, cx + Cm(0.5), cy + Cm(0.3), Cm(8), Cm(0.5),
                t, fs=14, color=PURPLE, bold=True)
    add_textbox(slide, cx + Cm(0.5), cy + Cm(1), Cm(9), Cm(1.8),
                d, fs=10, color=GRAY, ls=1.6)

# 底部路线
add_textbox(slide, Cm(2), Cm(10), Cm(30), Cm(0.5),
            "Phase 1 MVP(已完成) → Phase 2 商业化启动(近期) → Phase 3 平台化(中期)",
            fs=12, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

add_footer(slide, 14)


# ============================================================
# Slide 15: 结尾（修改模板最后一页）
# ============================================================
slide = prs.slides[len(prs.slides) - 1]
# 模板最后页已有THANKS，添加副标题
add_textbox(slide, Cm(5), Cm(10), Cm(25), Cm(0.8),
            "青囊AI · 以AI之眼观舌象 · 以古人之智辨体质", fs=14, color=GRAY,
            align=PP_ALIGN.CENTER)
add_textbox(slide, Cm(5), Cm(11.5), Cm(25), Cm(0.6),
            "访问地址：http://127.0.0.1:5000", fs=11, color=GRAY,
            align=PP_ALIGN.CENTER)


# ============================================================
# 保存
# ============================================================
prs.save(OUTPUT)
print(f"PPT已保存至：{OUTPUT}")
print(f"共 {len(prs.slides)} 页幻灯片")
